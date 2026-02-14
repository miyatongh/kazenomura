#!/usr/bin/env python3
"""
2026/2/16 Phase 0 コンサルティング説明会スライド生成スクリプト
全33枚のプレゼンテーションを python-pptx で作成（仕様書3.3/3.4追加版）
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ── 定数 ──
SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

# カラーパレット（落ち着いたトーン）
C_NAVY = RGBColor(0x1B, 0x2A, 0x4A)       # メインタイトル
C_DARK = RGBColor(0x2C, 0x3E, 0x50)       # 本文
C_BLUE = RGBColor(0x2E, 0x86, 0xC1)       # アクセント青
C_TEAL = RGBColor(0x17, 0xA5, 0x89)       # アクセント緑
C_ORANGE = RGBColor(0xE6, 0x7E, 0x22)     # 強調オレンジ
C_LIGHT_BG = RGBColor(0xF0, 0xF4, 0xF8)   # 薄い背景
C_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
C_GRAY = RGBColor(0x7F, 0x8C, 0x8D)       # サブテキスト
C_SECTION_BG = RGBColor(0x1B, 0x2A, 0x4A) # セクション扉背景

prs = Presentation()
prs.slide_width = SLIDE_WIDTH
prs.slide_height = SLIDE_HEIGHT


# ── ヘルパー関数 ──

def add_bg_shape(slide, color):
    """スライド全面に背景色の矩形を追加"""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Emu(0), Emu(0), SLIDE_WIDTH, SLIDE_HEIGHT
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_textbox(slide, left, top, width, height, text, font_size=18,
                bold=False, color=C_DARK, alignment=PP_ALIGN.LEFT,
                font_name="Yu Gothic"):
    """テキストボックスを追加"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_multiline_textbox(slide, left, top, width, height, lines,
                          font_size=18, color=C_DARK, line_spacing=1.5,
                          font_name="Yu Gothic", bold=False):
    """複数行テキストボックスを追加（各行を別パラグラフとして）"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line_data in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()

        if isinstance(line_data, dict):
            p.text = line_data.get("text", "")
            p.font.size = Pt(line_data.get("size", font_size))
            p.font.bold = line_data.get("bold", bold)
            p.font.color.rgb = line_data.get("color", color)
            p.font.name = line_data.get("font", font_name)
            p.alignment = line_data.get("align", PP_ALIGN.LEFT)
        else:
            p.text = str(line_data)
            p.font.size = Pt(font_size)
            p.font.bold = bold
            p.font.color.rgb = color
            p.font.name = font_name
        p.space_after = Pt(font_size * (line_spacing - 1))
    return txBox


def add_page_number(slide, num, total=33, color=C_GRAY):
    """右下にページ番号を追加"""
    add_textbox(slide, Inches(11.8), Inches(6.9), Inches(1.2), Inches(0.4),
                f"{num} / {total}", font_size=10, color=color,
                alignment=PP_ALIGN.RIGHT)


def make_title_slide(title, subtitle, date_info, presenter):
    """表紙スライド"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    add_bg_shape(slide, C_NAVY)

    # 上部ライン
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(1), Inches(1.8), Inches(11.3), Pt(3)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = C_TEAL
    line.line.fill.background()

    add_textbox(slide, Inches(1), Inches(2.0), Inches(11.3), Inches(1.5),
                title, font_size=36, bold=True, color=C_WHITE,
                alignment=PP_ALIGN.LEFT)
    add_textbox(slide, Inches(1), Inches(3.5), Inches(11.3), Inches(0.8),
                subtitle, font_size=22, color=RGBColor(0xAE, 0xBF, 0xD5),
                alignment=PP_ALIGN.LEFT)

    # 下部ライン
    line2 = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(1), Inches(5.2), Inches(11.3), Pt(1)
    )
    line2.fill.solid()
    line2.fill.fore_color.rgb = C_TEAL
    line2.line.fill.background()

    add_textbox(slide, Inches(1), Inches(5.4), Inches(5), Inches(0.5),
                date_info, font_size=16, color=RGBColor(0xAE, 0xBF, 0xD5))
    add_textbox(slide, Inches(7), Inches(5.4), Inches(5.3), Inches(0.5),
                presenter, font_size=16, color=RGBColor(0xAE, 0xBF, 0xD5),
                alignment=PP_ALIGN.RIGHT)
    return slide


def make_section_slide(number, title, subtitle=""):
    """セクション扉スライド"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg_shape(slide, C_SECTION_BG)

    # 大きな番号
    add_textbox(slide, Inches(1), Inches(1.5), Inches(2), Inches(2.5),
                str(number), font_size=96, bold=True,
                color=RGBColor(0x2E, 0x86, 0xC1))

    # タイトル
    add_textbox(slide, Inches(1), Inches(3.5), Inches(11), Inches(1.2),
                title, font_size=36, bold=True, color=C_WHITE)

    if subtitle:
        add_textbox(slide, Inches(1), Inches(4.8), Inches(11), Inches(0.8),
                    subtitle, font_size=20,
                    color=RGBColor(0xAE, 0xBF, 0xD5))
    return slide


def make_content_slide(slide_num, title, bullets, note="", emphasis_indices=None):
    """標準コンテンツスライド"""
    if emphasis_indices is None:
        emphasis_indices = []

    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # ヘッダー帯
    header_bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Emu(0), Emu(0), SLIDE_WIDTH, Inches(1.2)
    )
    header_bg.fill.solid()
    header_bg.fill.fore_color.rgb = C_NAVY
    header_bg.line.fill.background()

    add_textbox(slide, Inches(0.8), Inches(0.15), Inches(11.5), Inches(0.9),
                title, font_size=28, bold=True, color=C_WHITE)

    # 本文
    txBox = slide.shapes.add_textbox(
        Inches(1.0), Inches(1.6), Inches(11.3), Inches(5.0)
    )
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()

        if isinstance(bullet, dict):
            p.text = bullet.get("text", "")
            p.font.size = Pt(bullet.get("size", 20))
            p.font.bold = bullet.get("bold", i in emphasis_indices)
            p.font.color.rgb = bullet.get("color", C_ORANGE if i in emphasis_indices else C_DARK)
            p.font.name = bullet.get("font", "Yu Gothic")
        else:
            p.text = str(bullet)
            p.font.size = Pt(20)
            p.font.bold = (i in emphasis_indices)
            p.font.color.rgb = C_ORANGE if i in emphasis_indices else C_DARK
            p.font.name = "Yu Gothic"
        p.space_after = Pt(8)

    # ノート
    if note:
        add_textbox(slide, Inches(1.0), Inches(6.5), Inches(11.3), Inches(0.5),
                    note, font_size=14, color=C_GRAY)

    add_page_number(slide, slide_num)
    return slide


def make_two_column_slide(slide_num, title, left_title, left_items,
                          right_title, right_items, note=""):
    """2カラムスライド"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # ヘッダー帯
    header_bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Emu(0), Emu(0), SLIDE_WIDTH, Inches(1.2)
    )
    header_bg.fill.solid()
    header_bg.fill.fore_color.rgb = C_NAVY
    header_bg.line.fill.background()

    add_textbox(slide, Inches(0.8), Inches(0.15), Inches(11.5), Inches(0.9),
                title, font_size=28, bold=True, color=C_WHITE)

    # 左カラム背景
    left_bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(1.5),
        Inches(5.6), Inches(5.2)
    )
    left_bg.fill.solid()
    left_bg.fill.fore_color.rgb = RGBColor(0xEB, 0xEF, 0xF5)
    left_bg.line.fill.background()

    add_textbox(slide, Inches(0.8), Inches(1.6), Inches(5.2), Inches(0.6),
                left_title, font_size=22, bold=True, color=C_BLUE)

    add_multiline_textbox(slide, Inches(0.8), Inches(2.3), Inches(5.2), Inches(4.0),
                          left_items, font_size=18, color=C_DARK, line_spacing=1.6)

    # 右カラム背景
    right_bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(6.8), Inches(1.5),
        Inches(5.8), Inches(5.2)
    )
    right_bg.fill.solid()
    right_bg.fill.fore_color.rgb = RGBColor(0xE8, 0xF8, 0xF5)
    right_bg.line.fill.background()

    add_textbox(slide, Inches(7.0), Inches(1.6), Inches(5.4), Inches(0.6),
                right_title, font_size=22, bold=True, color=C_TEAL)

    add_multiline_textbox(slide, Inches(7.0), Inches(2.3), Inches(5.4), Inches(4.0),
                          right_items, font_size=18, color=C_DARK, line_spacing=1.6)

    if note:
        add_textbox(slide, Inches(1.0), Inches(6.8), Inches(11.3), Inches(0.5),
                    note, font_size=14, color=C_GRAY)

    add_page_number(slide, slide_num)
    return slide


# ========================================================================
# スライド生成
# ========================================================================

# ── Slide 1: 表紙 ──
make_title_slide(
    title="Phase 0「要件設計」コンサルティングのご説明",
    subtitle="何を、なぜ、どのように進めるか",
    date_info="2026年2月16日",
    presenter="株式会社PreSoft / eMu　上田昌夫"
)

# ── Slide 2: 本日のお話の位置づけ ──
make_two_column_slide(
    slide_num=2,
    title="本日のお話の位置づけ ── これまでとの違い",
    left_title="これまで（12月・1月）",
    left_items=[
        "■ 現状の構造的課題の可視化",
        "　 → 3つの構造的要因を特定",
        "",
        "■ 「仕事ナビ」構想の提言",
        "　 → 3原則と改革の方向性",
        "",
        "■ 変革の方向性への合意",
    ],
    right_title="本日（2月）",
    right_items=[
        "■ コンサルティングの進め方の説明",
        "　 → なぜこの方法なのか",
        "",
        "■ Phase 0 で何をするか",
        "　 → 具体的な活動と成果物",
        "",
        "■ 皆さまにお願いすること",
    ],
    note="今日は「何を・どう進めるか」のお話です"
)

# ── Slide 3: 前回までの確認 ──
make_content_slide(
    slide_num=3,
    title="確認：私たちが共有した「現状認識」",
    bullets=[
        "ヒアリングから、3つの構造的課題が明らかになりました：",
        "",
        "① データの「つなぎ目」の分断",
        "　  → 多重入力、月次決算が翌月25日まで確定しない",
        "",
        "② ルール・知識の属人化",
        "　  → 按分ルール、手当計算、請求条件が担当者の頭の中に",
        "",
        "③ 業務の不可視と「依頼-実行」関係の不全",
        "　  → 作業量が見えない、問い合わせが1日30-50件",
    ],
    note="これは「個人の問題」ではなく「仕組みの問題」です（12月報告書にて共有済み）",
    emphasis_indices=[2, 5, 8]
)

# ── Slide 4: 合意いただいた方向性 ──
make_content_slide(
    slide_num=4,
    title="確認：合意いただいた改革の方向性",
    bullets=[
        "3つの変革原則：",
        "　❶ 源流入力による一気通貫（リアルタイム・一個流し）",
        "　❷ 「知識」と「作業」の分離（ナレッジ・セントラル）",
        "　❸ 計画と実績による業務の可視化（PDマネジメント）",
        "",
        "実現のロードマップ：",
        {"text": "　Phase 0：都市計画の策定（3ヶ月）← 今日はここの中身です", "size": 22, "bold": True, "color": C_ORANGE},
        "　Phase 1：業務構造の整理と「止血」",
        "　Phase 2：パイロット検証",
        "　Phase 3：全社展開と継続的進化",
    ],
    emphasis_indices=[1, 2, 3]
)

# ── Slide 5: 本日の4つのテーマ ──
make_content_slide(
    slide_num=5,
    title="本日ご説明する4つのテーマ",
    bullets=[
        "",
        {"text": "①  なぜこの進め方なのか", "size": 24, "bold": True, "color": C_BLUE},
        "　　── コンサルティングの基本的な考え方",
        "",
        {"text": "②  何を・何のために行うか", "size": 24, "bold": True, "color": C_BLUE},
        "　　── Phase 0 の具体的な活動内容",
        "",
        {"text": "③  何が出来上がるのか", "size": 24, "bold": True, "color": C_BLUE},
        "　　── 成果物のイメージ",
        "",
        {"text": "④  皆さまにお願いすること", "size": 24, "bold": True, "color": C_BLUE},
        "　　── 風の村としてのご協力事項",
    ],
    note="ご説明 約40分 ＋ 質疑応答 約20分"
)

# ── Slide 6: セクション扉① ──
make_section_slide("①", "なぜこの進め方なのか", "コンサルティングの基本的な考え方")

# ── Slide 7: よくある失敗パターン ──
make_content_slide(
    slide_num=7,
    title="よくある失敗パターン（風の村でも起きていること）",
    bullets=[
        "「困った → ベンダーに相談 → パッケージ導入 → 合わない → 手作業が残る」",
        "",
        "■ 風の村での実例：",
        "　・「ほのぼの」── 請求書の名寄せに非対応 → 紙5,000通の発送が残った",
        "　・ 給与システム ── 特殊手当を計算できず、手で上書き（強制入力）",
        "　・ 按分処理 ── 「楽々精算」等で対応しきれず手作業が継続",
        "　・ ベンダーに相談 →「膨大なコストがかかるから今のままの方がいい」",
        "",
        {"text": "なぜこうなるか：業務の設計図なしにシステムを選んでいるから", "size": 22, "bold": True, "color": C_ORANGE},
    ],
    note="比喩：「既製服を買ってから体に合わせて切る」のではなく「まず採寸してオーダーメイドする」"
)

# ── Slide 8: 都市計画 ──
make_two_column_slide(
    slide_num=8,
    title="私たちの考え方：「ビル建て替え」ではなく「都市計画」",
    left_title="ビル建て替え（従来型）",
    left_items=[
        "・古いシステムを新しく置き換える",
        "・一発勝負（失敗すると大損害）",
        "・完成するまで効果が見えない",
        "・終わったら「次の建て替え」まで放置",
    ],
    right_title="都市計画（私たちのアプローチ）",
    right_items=[
        "・まずマスタープランを描く",
        "・区画ごとに段階的に整備",
        "・今の建物を使いながら進める",
        "・終わりのない「まちづくり」",
    ],
    note="Phase 0 ＝ マスタープラン策定。「設計図なしに工事は始めません」"
)

# ── Slide 9: 4つの視点 ──
make_content_slide(
    slide_num=9,
    title="業務を「4つの視点」で設計する",
    bullets=[
        "都市計画が道路・水道・用途地域・建築基準をセットで設計するように、",
        "業務改革も4つの視点をセットで設計します：",
        "",
        {"text": "❶ 業務の流れ", "size": 22, "bold": True, "color": C_BLUE},
        "　　誰が何をどの順番で行うか（＝道路計画）",
        "",
        {"text": "❷ 情報の流れ", "size": 22, "bold": True, "color": C_BLUE},
        "　　どのデータがどこで生まれ、どこに届くか（＝上下水道計画）",
        "",
        {"text": "❸ 道具の配置", "size": 22, "bold": True, "color": C_BLUE},
        "　　どのシステムがどの仕事を助けるか（＝用途地域計画）",
        "",
        {"text": "❹ 基盤", "size": 22, "bold": True, "color": C_BLUE},
        "　　全体を支えるインフラ（＝建築基準）",
    ],
    note="この4つをバラバラに検討すると「つなぎ目」が生まれます。だからセットで設計します"
)

# ── Slide 10: 【新規】組織を「3つの階層」で捉える ──
make_content_slide(
    slide_num=10,
    title="組織を「3つの階層」で捉える",
    bullets=[
        "風の村の活動全体を一つのシステムとして見ると、3つの階層に整理できます：",
        "",
        {"text": "❶ サービス提供の層", "size": 22, "bold": True, "color": C_BLUE},
        "　　利用者からの依頼を受け、サービスを提供し、記録・請求する流れ",
        "",
        {"text": "❷ 経営資源を管理する層", "size": 22, "bold": True, "color": C_BLUE},
        "　　人・モノ・お金・時間を配分し管理する仕事（人事給与、経理等）",
        "",
        {"text": "❸ 現場で実行する層", "size": 22, "bold": True, "color": C_BLUE},
        "　　日々のシフト実行、入金処理、物品発注など実世界との接点",
        "",
        {"text": "各階層の中に「知識（ルール）」「計画と実績」「報告」の3つの側面があります", "size": 18, "bold": False, "color": C_GRAY},
    ],
    note="今、風の村では層と層の「つなぎ目」で情報が途切れている。これが多重入力の根本原因です"
)

# ── Slide 11: 【新規】仕事の流れを設計する3つの原則 ──
make_content_slide(
    slide_num=11,
    title="仕事の流れを設計する3つの原則",
    bullets=[
        "私たちが業務の流れを設計するときの3つの原則：",
        "",
        {"text": "❶ 計画 → 実行 → 確認（PDマネジメント）", "size": 22, "bold": True, "color": C_BLUE},
        "　　計画を立て、実績を記録し、差を見て修正する",
        "　　計画通りにいかないことが前提。リアルタイムに状況を把握する",
        "",
        {"text": "❷ 依頼 → 約束 → 実行 → 検収", "size": 22, "bold": True, "color": C_BLUE},
        "　　すべての仕事は「頼む → 引き受ける → やる → 確認する」の4ステップ",
        "　　比喩：レストランの注文 → 厨房が受ける → 調理 → お客様が確認",
        "",
        {"text": "❸ リアルタイム・一つずつ確実に", "size": 22, "bold": True, "color": C_BLUE},
        "　　月末にまとめて処理するのではなく、発生のたびに完結させる",
    ],
    note="この原則がシステムに反映されていないと、途中で情報が途切れます"
)

# ── Slide 12: 【新規】この原則を当てはめると何が変わるか ──
make_two_column_slide(
    slide_num=12,
    title="この原則を当てはめると何が変わるか",
    left_title="現状",
    left_items=[
        "・月末にまとめて按分計算",
        "　→ 月末の業務集中と遅延",
        "",
        "・紙の請求書を手作業で照合",
        "　→ 5,000通の紙が毎月発生",
        "",
        "・シフト確定が25日ギリギリ",
        "　→ 人件費の見通しが立たない",
        "",
        "・ベンダーに「変えられない」と言われる",
    ],
    right_title="原則適用後",
    right_items=[
        "・サービス記録時点で按分が自動連携",
        "　→ リアルタイムで経営数値を把握",
        "",
        "・請求データが日次で確定",
        "　→ 紙を大幅に削減",
        "",
        "・計画と実績の差がリアルタイムで見える",
        "　→ 早期の軌道修正が可能",
        "",
        "・設計図があるから要望を正確に伝えられる",
    ],
    note="原則に沿って業務を再設計するから、システムが変わった後も仕組みが崩れません"
)

# ── Slide 13: ヒアリングの方法 ──
make_content_slide(
    slide_num=13,
    title="ヒアリングの方法：「聴く」から「構造化する」へ",
    bullets=[
        "■ 一般的なヒアリング：",
        "　「何に困っていますか？」→ バラバラの要望リスト → 構造がない",
        "",
        {"text": "■ 私たちのヒアリング：", "size": 20, "bold": True, "color": C_BLUE},
        "　 散らばった声から「共通の型（パターン）」を見つける",
        "",
        "　・質的研究手法（M-GTA等）：一人ひとりの話から組織全体の型を発見",
        "　・KJ法：バラバラの情報を構造的に整理",
        "　・生成AI：大量の聞き取り情報を高速にパターン化する道具",
        "",
        {"text": "11月のヒアリングで「個別の愚痴」が「3つの共通構造問題」に集約できたのは、", "size": 18, "bold": False, "color": C_GRAY},
        {"text": "この手法を使ったからです。", "size": 18, "bold": False, "color": C_GRAY},
    ]
)

# ── Slide 14: 要件設計とは ──
make_content_slide(
    slide_num=14,
    title="「要件設計」とは何か",
    bullets=[
        {"text": "「要件設計」はシステムの仕様書ではありません。", "size": 22, "bold": True, "color": C_ORANGE},
        "",
        "「風の村の業務が、将来どのように流れるべきか」を描いた設計図です。",
        "",
        "4つの階層で設計します：",
        "　❶ 業務要件　── 「こういう仕事の流れにしたい」",
        "　❷ 運用要件　── 「現場ではこう使いたい」",
        "　❸ 機能要件　── 「こういうことができる道具がほしい」",
        "　❹ システム要件 ── 「その道具を動かすための条件」",
    ],
    note="比喩：注文住宅を建てるとき、施主が「どう暮らしたいか」を伝え、建築士が「間取り図」と「設備仕様書」に翻訳する"
)

# ── Slide 15: 3つの約束 ──
make_content_slide(
    slide_num=15,
    title="私たちの3つの約束",
    bullets=[
        "",
        "",
        {"text": "❶  全体を見てから、部分を決めます", "size": 28, "bold": True, "color": C_NAVY},
        "　　 部分最適ではなく、全体最適を設計する",
        "",
        "",
        {"text": "❷  現場の言葉を、設計の言葉に翻訳します", "size": 28, "bold": True, "color": C_NAVY},
        "　　 技術ではなく、業務の視点から入る",
        "",
        "",
        {"text": "❸  設計図ができてから、工事を始めます", "size": 28, "bold": True, "color": C_NAVY},
        "　　 手戻りを最小化する",
    ]
)

# ── Slide 16: セクション扉② ──
make_section_slide("②", "何を・何のために行うか", "Phase 0「要件設計」の具体的な進め方")

# ── Slide 17: Phase 0 全体像 ──
make_content_slide(
    slide_num=17,
    title="Phase 0 の全体像：3ヶ月間の活動",
    bullets=[
        {"text": "3つの活動を並行して進めます：", "size": 22, "bold": True, "color": C_NAVY},
        "",
        {"text": "A. ヒアリング調査（約2ヶ月間）", "size": 22, "bold": True, "color": C_BLUE},
        "　　 1日2回 × 週2日 × 約2ヶ月 ＝ 25～30セッション",
        "",
        {"text": "B. 隔週プロジェクト会議（全期間）", "size": 22, "bold": True, "color": C_BLUE},
        "　　 2週間に1回 × 約3ヶ月 ＝ 6回程度",
        "",
        {"text": "C. 分析・設計作業（コンサルタント側）", "size": 22, "bold": True, "color": C_BLUE},
        "　　 ヒアリング → 構造化 → 設計 → マスタープラン策定",
    ],
    note="3ヶ月後のアウトプット：マスタープラン ＋ 要件定義書 ＋ マスターデータ構造設計"
)

# ── Slide 18: ヒアリング調査 ──
make_content_slide(
    slide_num=18,
    title="活動① ヒアリング調査 ── 皆さまの「仕事」を丸ごと理解する",
    bullets=[
        "■ 規模：約25～30回のセッション（1回 90分目安）",
        "■ ペース：1日2回 × 週2日 × 約2ヶ月",
        "■ 対象：管理本部の各部署 ＋ 現場（事業所）",
        "",
        "■ ヒアリングで聞くこと：",
        "　・今のお仕事の流れを、最初から最後まで教えてください",
        "　・困っていること、工夫していること、変えたいこと",
        "　・他の部署や現場とのやりとり",
        "",
        {"text": "「技術の話はしません。仕事の話を聞かせてください」", "size": 22, "bold": True, "color": C_TEAL},
    ],
    note="11月に管理本部5部署で実施したヒアリングを、全部署・現場に広げるイメージです"
)

# ── Slide 19: プロジェクト会議 ──
make_content_slide(
    slide_num=19,
    title="活動② 隔週プロジェクト会議 ── 経過の共有と方向性の確認",
    bullets=[
        "■ 頻度：2週間に1回（約1時間）",
        "■ 参加：風の村側プロジェクト責任者・主要メンバー ＋ コンサルタント",
        "",
        "■ 各回の内容：",
        "　・ヒアリングで見えてきたことの中間報告",
        "　・方向性の確認と軌道修正",
        "　・次の2週間の計画共有",
        "",
        {"text": "「密室で分析するのではなく、進捗を常に共有します」", "size": 22, "bold": True, "color": C_TEAL},
    ],
    note="進め方に疑問や懸念があれば、この場でいつでも方向修正できます"
)

# ── Slide 20: 裏側の作業 ──
make_content_slide(
    slide_num=20,
    title="活動③ 分析・設計作業 ── ヒアリングの「裏側」",
    bullets=[
        "コンサルタント側で、ヒアリングの合間に以下の作業を進めます：",
        "",
        "　Step 1  ヒアリング内容の文字起こし・整理（AI活用）",
        "　Step 2  業務フローの可視化（誰が・何を・いつ・どのデータで）",
        "　Step 3  課題の分類と構造化（パターンの抽出）",
        "　Step 4  あるべき業務モデルの設計",
        "　Step 5  システム要件への変換",
        "",
        {"text": "1回1回のヒアリングが、設計図の一部になっていきます", "size": 22, "bold": True, "color": C_TEAL},
    ]
)

# ── Slide 21: なぜ25-30回 ──
make_content_slide(
    slide_num=21,
    title="なぜ25～30回もの聞き取りが必要なのか",
    bullets=[
        "",
        {"text": "❶ 約50拠点の多様な業務を網羅するため", "size": 22, "bold": True, "color": C_NAVY},
        "　　介護・保育・障害の各事業で業務が異なる",
        "",
        {"text": "❷ 部署間の「つなぎ目」を見つけるため", "size": 22, "bold": True, "color": C_NAVY},
        "　　構造的な問題は、部署と部署の接点に潜んでいる",
        "",
        {"text": "❸ 十分な事例数で「共通の型」を抽出するため", "size": 22, "bold": True, "color": C_NAVY},
        "　　少数の聞き取りでは、偏った設計図になってしまう",
        "",
        {"text": "❹ 設計の精度を上げ、後の手戻りを防ぐため", "size": 22, "bold": True, "color": C_NAVY},
        "　　設計図が粗いと、工事段階で想定外の問題が起きる",
    ],
    note="比喩：健康診断で血液検査1項目だけでは診断できない。全身を調べるから的確な治療方針が立てられる"
)

# ── Slide 22: 都市計画に例えると ──
make_content_slide(
    slide_num=22,
    title="都市計画に例えると ── Phase 0 の位置づけ",
    bullets=[
        "",
        {"text": "ヒアリング ＝ 住民への聞き取り調査", "size": 24, "bold": True, "color": C_BLUE},
        "　今どこに住んで、どう暮らしているか",
        "",
        {"text": "業務分析 ＝ 地質調査・交通量調査", "size": 24, "bold": True, "color": C_BLUE},
        "　地盤の強さや人の流れを把握する",
        "",
        {"text": "要件設計 ＝ マスタープラン策定", "size": 24, "bold": True, "color": C_BLUE},
        "　ゾーニング、道路計画、優先整備エリアを決める",
        "",
        "",
        {"text": "Phase 0 では「工事」はしません。設計図を共有し、合意してから着工します。", "size": 20, "bold": True, "color": C_ORANGE},
    ]
)

# ── Slide 23: セクション扉③ ──
make_section_slide("③", "何が出来上がるのか", "Phase 0 の成果物")

# ── Slide 24: 成果物の全体像 ──
make_content_slide(
    slide_num=24,
    title="Phase 0 で皆さまにお渡しするもの",
    bullets=[
        "",
        {"text": "3つの成果物：", "size": 24, "bold": True, "color": C_NAVY},
        "",
        {"text": "❶  マスタープラン（基本計画 ＋ 実施計画）", "size": 24, "bold": True, "color": C_BLUE},
        "　　「どこから・どう変えるか」の全体計画",
        "",
        {"text": "❷  要件定義書", "size": 24, "bold": True, "color": C_BLUE},
        "　　「仕事ナビ」に何が必要かの仕様",
        "",
        {"text": "❸  マスターデータ構造設計", "size": 24, "bold": True, "color": C_BLUE},
        "　　「情報の背骨」の設計図",
    ]
)

# ── Slide 25: マスタープラン ──
make_content_slide(
    slide_num=25,
    title="① マスタープラン ── 「どこから・どう変えるか」の全体計画",
    bullets=[
        {"text": "■ 基本計画", "size": 22, "bold": True, "color": C_BLUE},
        "　・現状の業務構造の全体像（可視化された業務フロー）",
        "　・あるべき姿の設計（仕事ナビ構想に基づく将来モデル）",
        "　・ギャップ分析（現状と将来像のギャップ一覧）",
        "",
        {"text": "■ 実施計画", "size": 22, "bold": True, "color": C_BLUE},
        "　・優先順位の提案（どの業務エリアから着手するか）",
        "　・フェーズ分けと概算スケジュール",
        "　・必要なリソースと投資規模の概算",
        "",
        {"text": "経営判断の材料として使えるものを作ります", "size": 22, "bold": True, "color": C_ORANGE},
    ]
)

# ── Slide 26: 要件定義書 ──
make_content_slide(
    slide_num=26,
    title="② 要件定義書 ── 「仕事ナビ」に何が必要かの仕様",
    bullets=[
        "4つの階層で、業務の「設計図」を描きます：",
        "",
        "　❶ 業務要件　── 新しい業務の流れの定義",
        "　　（例：請求処理はこう変わる）",
        "",
        "　❷ 運用要件　── 日常の運用ルールの定義",
        "　　（例：データ入力の締め日はいつか）",
        "",
        "　❸ 機能要件　── システムに必要な機能の一覧",
        "　　（例：自動按分計算、リアルタイム進捗表示）",
        "",
        "　❹ システム要件 ── 技術面の条件",
    ],
    note="比喩：注文住宅の「間取り図」と「設備仕様書」。これがあれば、ベンダーに要望を正しく伝えられます"
)

# ── Slide 27: マスターデータ構造設計 ──
make_two_column_slide(
    slide_num=27,
    title="③ マスターデータ構造設計 ── 「情報の背骨」の設計図",
    left_title="現状（バラバラ）",
    left_items=[
        "・ほのぼの、MJS、弥生、勤怠システム…",
        "　 それぞれに職員情報・事業所情報を登録",
        "",
        "・同じ職員が各システムで別の番号",
        "・事業所の名称表記がシステムごとに違う",
        "",
        "→ これが多重入力と不整合の根本原因",
    ],
    right_title="あるべき姿（統一）",
    right_items=[
        "・全システム共通の「元データ」を設計",
        "",
        "・職員マスタ、事業所マスタ、",
        "　 勘定科目マスタ、サービスマスタ等",
        "",
        "・一箇所で更新すれば全体に反映",
        "",
        "→ 「One Fact, One Place」の実現",
    ],
    note="比喩：住所体系の統一。番地がバラバラでは郵便も届きません"
)

# ── Slide 28: セクション扉④ ──
make_section_slide("④", "皆さまにお願いすること", "風の村としてのご協力事項")

# ── Slide 29: 4つのお願い ──
make_content_slide(
    slide_num=29,
    title="Phase 0 を成功させるための4つのお願い",
    bullets=[
        "",
        {"text": "❶  ヒアリング対象者の選定と日程調整", "size": 24, "bold": True, "color": C_NAVY},
        "　　各部署・事業所から、実務を担当されている方をご指名ください",
        "",
        {"text": "❷  既存資料・システム情報のご提供", "size": 24, "bold": True, "color": C_NAVY},
        "　　マニュアル、帳票、システム仕様書等（完璧でなくて構いません）",
        "",
        {"text": "❸  隔週プロジェクト会議へのご参加", "size": 24, "bold": True, "color": C_NAVY},
        "　　2週間に1回、約1時間の経過共有の場です",
        "",
        {"text": "❹  プロジェクト責任者・連絡窓口のご決定", "size": 24, "bold": True, "color": C_NAVY},
        "　　意思決定者と日常のやりとり担当者をそれぞれ1名",
    ]
)

# ── Slide 30: なぜ一緒にか ──
make_content_slide(
    slide_num=30,
    title="なぜ「お任せ」ではなく「一緒に」なのか",
    bullets=[
        "",
        "コンサルタントは、業務設計の専門家です。",
        "",
        {"text": "しかし、風の村の業務の専門家は、皆さまです。", "size": 24, "bold": True, "color": C_NAVY},
        "",
        "設計図は、使う人と一緒に描かなければ、",
        "使える設計図になりません。",
        "",
        "一緒に作った設計図だからこそ、",
        "実行段階で「自分たちのもの」になります。",
        "",
        {"text": "大きな体制は不要です。少人数で機動的に進めます。", "size": 20, "bold": False, "color": C_GRAY},
    ],
    note="比喩：注文住宅は、打ち合わせに参加した人ほど完成後の満足度が高い"
)

# ── Slide 31: ヒアリングを受ける方へ ──
make_content_slide(
    slide_num=31,
    title="ヒアリングを受ける方へのメッセージ",
    bullets=[
        "",
        "",
        {"text": "正解を答える場ではありません。", "size": 28, "bold": True, "color": C_NAVY},
        "",
        {"text": "日頃の仕事を、そのまま教えてください。", "size": 28, "bold": True, "color": C_NAVY},
        "",
        "",
        "困っていること、工夫していること、",
        "「なんでこうなってるんだろう」と感じていること。",
        "",
        "何を言っても大丈夫です。誰が何を言ったかは報告しません。",
    ],
    note="部長の皆さまから現場の方への声掛けをお願いします。「自由に話していい」という安心感がヒアリングの質を決めます"
)

# ── Slide 32: まとめ ──
slide29 = prs.slides.add_slide(prs.slide_layouts[6])
header_bg = slide29.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, Emu(0), Emu(0), SLIDE_WIDTH, Inches(1.2)
)
header_bg.fill.solid()
header_bg.fill.fore_color.rgb = C_NAVY
header_bg.line.fill.background()
add_textbox(slide29, Inches(0.8), Inches(0.15), Inches(11.5), Inches(0.9),
            "まとめ：Phase 0 で実現すること", font_size=28, bold=True, color=C_WHITE)

# 4象限
labels = [
    ("考え方", "都市計画アプローチ\n4つの視点 × 3つの階層\n3つの原則で業務プロセス設計"),
    ("活動内容", "ヒアリング25-30回\n隔週プロジェクト会議\n分析・設計作業（3ヶ月）"),
    ("成果物", "マスタープラン\n要件定義書\nマスターデータ構造設計"),
    ("お願い", "対象者選定・日程調整\n資料提供\n会議参加・体制決定"),
]
colors = [C_BLUE, C_TEAL, C_BLUE, C_TEAL]
positions = [
    (Inches(0.8), Inches(1.5)),
    (Inches(6.8), Inches(1.5)),
    (Inches(0.8), Inches(4.3)),
    (Inches(6.8), Inches(4.3)),
]

for idx, ((label, content), color, (x, y)) in enumerate(zip(labels, colors, positions)):
    box = slide29.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(5.5), Inches(2.5)
    )
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(0xF7, 0xF9, 0xFC) if idx % 2 == 0 else RGBColor(0xEA, 0xF7, 0xF4)
    box.line.color.rgb = color

    num = f"①②③④"[idx]
    add_textbox(slide29, x + Inches(0.3), y + Inches(0.2), Inches(4.8), Inches(0.5),
                f"{num} {label}", font_size=20, bold=True, color=color)
    add_textbox(slide29, x + Inches(0.3), y + Inches(0.8), Inches(4.8), Inches(1.5),
                content, font_size=16, color=C_DARK)

add_page_number(slide29, 32)

# ── Slide 33: 次のステップ ──
make_content_slide(
    slide_num=33,
    title="次のステップ",
    bullets=[
        "",
        {"text": "☐  本日の内容に関するご質問・ご意見（この場で）", "size": 22, "bold": False, "color": C_DARK},
        "",
        {"text": "☐  プロジェクト責任者・連絡窓口のご指名", "size": 22, "bold": False, "color": C_DARK},
        "",
        {"text": "☐  ヒアリング対象部署・対象者リストの作成", "size": 22, "bold": False, "color": C_DARK},
        "",
        {"text": "☐  契約手続きの確認", "size": 22, "bold": False, "color": C_DARK},
        "",
        {"text": "☐  キックオフ会議の日程調整", "size": 22, "bold": False, "color": C_DARK},
    ],
    note="スケジュール案：3月キックオフ → 4月ヒアリング本格実施 → 5月分析・設計・最終報告"
)


# ── 保存 ──
output_path = "/Users/miyaton/Library/CloudStorage/GoogleDrive-ueda@miyaton.com/マイドライブ/eMu/@生活クラブ風の村/20260216_説明会/20260216_Phase0説明会.pptx"
prs.save(output_path)
print(f"✓ 保存完了: {output_path}")
print(f"  スライド数: {len(prs.slides)}")
